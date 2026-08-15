"""Assemble the finished slide images into a 16:9 PPTX (with speaker notes
attached to each slide) and a PDF.

Notes are parsed out of F11_final_script.md. Each "## Slide N ..." section
contributes two things to that slide's notes pane:
  1. the script — the blockquote lines you actually say, which are Arabic-dominant
  2. a DELIVERY section — the English guidance blocks and stage directions

They are separated so the pane opens on what you say, with the coaching beneath.
"""
import re
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

BASE = Path(r"F:\Desktop\graduation\research_decisions\defence")
SRC  = BASE / "slides" / "final_v3"
OUT  = BASE / "slides"
TMP  = Path(r"C:\Users\moham\AppData\Local\Temp\claude\f--Desktop-graduation\60f05023-192a-45f4-8e47-22a900f4a241\scratchpad\deckimg")
TMP.mkdir(exist_ok=True)

ORDER = ["slide_01_title","slide_02_agenda","slide_03_llm","slide_04_rag","slide_05_bottleneck",
"slide_06_literature","slide_07_arabic","slide_08_rq","slide_09_objectives","slide_10_system",
"slide_11_baseline","slide_12_query2doc","slide_13_densevsbm25","slide_14_repetition",
"slide_15_fivenouns","slide_16_csqe","slide_17_placement","slide_18_journey","slide_19_metrics",
"slide_20_conclusions","slide_21_futurework","slide_22_thanks"]

# ---------- parse the script ----------
AR = re.compile(r'[\u0600-\u06FF]')

def strip_md(s: str) -> str:
    s = re.sub(r'\*\*(.+?)\*\*', r'\1', s)
    s = re.sub(r'`(.+?)`', r'\1', s)
    s = re.sub(r'\[(.+?)\]\(.+?\)', r'\1', s)
    s = re.sub(r'^\s*>\s?', '', s)
    return s.rstrip()

def is_script(s: str) -> bool:
    """Script lines are Arabish and always contain Arabic. The English guidance
    blocks in F11 are pure Latin prose, so presence of any Arabic is the reliable
    signal -- 'Arabic-dominant' wrongly demoted lines heavy with English terms
    (e.g. slide 21: dialectal Arabic / first-pass quality gate / embedding models)."""
    return bool(AR.search(s))

md = (BASE / "F11_final_script.md").read_text(encoding="utf-8")

notes = {}          # slide number -> (title, [script lines], [delivery lines])
cur = None
block = []          # consecutive blockquote lines form one block

def flush(block, cur):
    if not block or cur is None:
        return
    text = "\n".join(block).strip()
    if not text:
        return
    title, script, delivery = notes[cur]
    (script if is_script(text) else delivery).append(text)

for raw in md.splitlines():
    line = raw.rstrip()
    m = re.match(r'^##\s+Slide\s+(\d+)\s*·\s*(.*)', line)
    if m:
        flush(block, cur); block = []
        cur = int(m.group(1))
        title = re.sub(r'\s*\*\(.*?\)\*\s*$', '', strip_md(m.group(2))).strip()
        notes[cur] = (title, [], [])
        continue
    if line.startswith('#'):
        flush(block, cur); block = []
        cur = None
        continue
    if cur is None:
        continue
    if line.startswith('>'):
        block.append(strip_md(line))
    elif line.startswith('*['):
        flush(block, cur); block = []
        notes[cur][2].append(strip_md(line).strip('*'))
    else:
        flush(block, cur); block = []
flush(block, cur)

# ---------- images ----------
W, H = 2560, 1440
small = []
for stem in ORDER:
    im = Image.open(SRC / f"{stem}.png").convert("RGB").resize((W, H), Image.LANCZOS)
    q = TMP / f"{stem}.jpg"
    im.save(q, "JPEG", quality=90, optimize=True)
    small.append((q, im))

# ---------- PPTX with notes ----------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

attached = 0
for i, (q, _) in enumerate(small, start=1):
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(str(q), 0, 0, width=prs.slide_width, height=prs.slide_height)

    title, script, delivery = notes.get(i, (ORDER[i-1], [], []))
    parts = [f"SLIDE {i} — {title}", ""]
    if script:
        parts += script + [""]
    if delivery:
        parts += ["— DELIVERY —"] + delivery
    s.notes_slide.notes_text_frame.text = "\n".join(parts).strip()
    if script:
        attached += 1

p1 = OUT / "Defence_Presentation.pptx"
prs.save(str(p1))
print(f"PPTX  {p1.stat().st_size/1024/1024:.1f} MB   notes on {attached}/22 slides")

# ---------- PDF ----------
pages = [im for _, im in small]
p2 = OUT / "Defence_Presentation.pdf"
pages[0].save(str(p2), save_all=True, append_images=pages[1:], resolution=150.0, quality=90)
print(f"PDF   {p2.stat().st_size/1024/1024:.1f} MB  {len(pages)} pages")
